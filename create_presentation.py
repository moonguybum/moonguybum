from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def create_presentation():
    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 슬라이드 마스터 설정 (제목과 내용이 없는 빈 화면 레이아웃)
    blank_slide_layout = prs.slide_layouts[6]

    # 공통 텍스트 상자 추가 함수
    def add_text_box(slide, text, left, top, width, height, font_size, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.CENTER):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox

    # [슬라이드 1] 절체절명의 순간 (도입)
    slide1 = prs.slides.add_slide(blank_slide_layout)
    # 배경을 어둡게 처리 (영상 삽입을 위한 안내)
    bg1 = slide1.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(30, 30, 30)
    add_text_box(slide1, "[영상 삽입] 어둡고 출렁이는 물결 (1인칭)", Inches(1), Inches(3), Inches(8), Inches(1), 24, color=RGBColor(200, 200, 200))

    # [슬라이드 2] 이성의 회복 (반전과 경각심)
    slide2 = prs.slides.add_slide(blank_slide_layout)
    bg2 = slide2.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 완전한 검은 배경 (암전 효과)
    add_text_box(slide2, "O₂ ↓", Inches(3), Inches(2.5), Inches(4), Inches(2), 80, bold=True, color=RGBColor(0, 176, 240))

    # [슬라이드 3] 골든타임의 액션 (전개 1)
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_text_box(slide3, "30 : 2", Inches(1), Inches(1.5), Inches(8), Inches(1.5), 70, bold=True, color=RGBColor(192, 0, 0))
    add_text_box(slide3, "모두 투명", Inches(1), Inches(3.5), Inches(8), Inches(1.5), 70, bold=True, color=RGBColor(255, 0, 0))
    add_text_box(slide3, "[배경 영상 삽입] 심폐소생술 흑백 클로즈업", Inches(1), Inches(5.5), Inches(8), Inches(1), 20, color=RGBColor(100, 100, 100))

    # [슬라이드 4] 또 다른 생명의 신호 (전환 브릿지)
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_text_box(slide4, "[이미지 삽입] 구출되는 작은 강아지 (흑백에서 컬러로)", Inches(1), Inches(3.5), Inches(8), Inches(1), 28, color=RGBColor(100, 100, 100))

    # [슬라이드 5] 반려견 CPR 기초 (전개 2)
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_text_box(slide5, "소형견 : 심장 바로 위", Inches(0.5), Inches(1.5), Inches(4.5), Inches(1), 32, bold=True)
    add_text_box(slide5, "중·대형견 : 가슴의 가장 높은 곳", Inches(5), Inches(1.5), Inches(4.5), Inches(1), 32, bold=True)
    add_text_box(slide5, "[이미지 삽입] 강아지 흉부 해부도 일러스트", Inches(2), Inches(4.5), Inches(6), Inches(1), 24, color=RGBColor(100, 100, 100))

    # [슬라이드 6] 인공호흡의 차이 (클라이맥스 강조)
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_text_box(slide6, "정반대 방향", Inches(1), Inches(2.5), Inches(8), Inches(2), 80, bold=True)
    add_text_box(slide6, "[영상 삽입] 코를 통해 바람을 불어넣는 3D 애니메이션", Inches(1), Inches(5), Inches(8), Inches(1), 24, color=RGBColor(100, 100, 100))

    # [슬라이드 7] 결론 및 행동 촉구 (감동적 마무리)
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_text_box(slide7, "생명의 신호에 답할 수 있는 사람", Inches(1), Inches(1.5), Inches(8), Inches(1), 36, bold=True)
    add_text_box(slide7, "90명", Inches(1), Inches(3), Inches(8), Inches(2), 100, bold=True)
    add_text_box(slide7, "[배경 이미지] 건강을 회복한 강아지와 주인의 따뜻한 사진", Inches(1), Inches(5.5), Inches(8), Inches(1), 20, color=RGBColor(100, 100, 100))

    # 파일 저장
    filename = "강아지가_보낸_구조신호_최종.pptx"
    prs.save(filename)
    print(f"PPTX 파일이 성공적으로 생성되었습니다: {filename}")
    print("이제 이 파일을 열어 안내된 위치에 영상과 이미지를 삽입하세요.")


if __name__ == "__main__":
    create_presentation()
